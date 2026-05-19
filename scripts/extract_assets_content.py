import json
import math
import re
import shutil
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.units import mm
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "UE LABS Presentation_2025.pptx"
ASSETS = ROOT / "assets"
EXTRACTED = ASSETS / "extracted"
PROCESSED = ASSETS / "processed"
REPORTS = ROOT / "reports"

NS = {
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

STANDARDS = [
    "IEC 62208",
    "ISO 12944-6",
    "ISO 4624",
    "IEC 60079-0",
    "EN 61386-1",
    "IEC 61439-1",
    "IEC 61439-5",
    "UL 50 / CSA C22.2 Nº 94.1",
    "UL 50E / CSA C22.2 Nº 94.2",
    "UL 746 C",
]

LAB_NAMES = [
    "UE Labs",
    "Capellades Lab",
    "Capellades",
    "Sarre-Union Lab",
    "Sarre-Union",
    "Molins de Rei Lab",
    "Molins de Rei",
    "Sarre Union",
]

PRODUCT_NAMES = ["Spacial", "Thalassa", "ClimaSys"]


def clean_text(value: str) -> str:
    value = value.replace("\x0b", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def rel_target_to_media(target: str) -> str | None:
    if not target:
        return None
    target = target.replace("\\", "/")
    if "media/" not in target:
        return None
    return Path(target).name


def collect_slide_media(zf: zipfile.ZipFile) -> dict[int, list[str]]:
    media_by_slide: dict[int, list[str]] = defaultdict(list)
    rel_paths = sorted(
        [p for p in zf.namelist() if p.startswith("ppt/slides/_rels/slide") and p.endswith(".xml.rels")],
        key=lambda p: int(re.search(r"slide(\d+)\.xml\.rels$", p).group(1)),
    )
    for rel_path in rel_paths:
        slide_num = int(re.search(r"slide(\d+)\.xml\.rels$", rel_path).group(1))
        root = ET.fromstring(zf.read(rel_path))
        for rel in root.findall("rel:Relationship", NS):
            rel_type = rel.attrib.get("Type", "")
            if "image" not in rel_type:
                continue
            media_name = rel_target_to_media(rel.attrib.get("Target", ""))
            if media_name and media_name not in media_by_slide[slide_num]:
                media_by_slide[slide_num].append(media_name)
    return media_by_slide


def collect_notes(zf: zipfile.ZipFile) -> dict[int, str]:
    notes_by_slide: dict[int, str] = {}
    slide_rels = sorted(
        [p for p in zf.namelist() if p.startswith("ppt/slides/_rels/slide") and p.endswith(".xml.rels")],
        key=lambda p: int(re.search(r"slide(\d+)\.xml\.rels$", p).group(1)),
    )
    for rel_path in slide_rels:
        slide_num = int(re.search(r"slide(\d+)\.xml\.rels$", rel_path).group(1))
        root = ET.fromstring(zf.read(rel_path))
        note_target = None
        for rel in root.findall("rel:Relationship", NS):
            if "notesSlide" in rel.attrib.get("Type", ""):
                note_target = rel.attrib.get("Target", "").replace("../", "ppt/")
                break
        if not note_target or note_target not in zf.namelist():
            notes_by_slide[slide_num] = ""
            continue
        note_root = ET.fromstring(zf.read(note_target))
        texts = [t.text or "" for t in note_root.findall(".//a:t", NS)]
        text = clean_text("\n".join([t for t in texts if t.strip()]))
        # Notes XML often carries slide number placeholders. Keep authored notes only.
        text = re.sub(r"^Slide \d+\s*", "", text).strip()
        notes_by_slide[slide_num] = text
    return notes_by_slide


def estimate_image_quality(width: int | None, height: int | None, file_type: str) -> str:
    if width is None or height is None:
        if file_type.lower() in {"emf", "svg"}:
            return "vector / not raster"
        return "unknown"
    shortest = min(width, height)
    area = width * height
    if shortest >= 1000 or area >= 1_250_000:
        return "high"
    if shortest >= 500 or area >= 400_000:
        return "medium"
    if shortest >= 220:
        return "low-medium"
    return "low"


def classify_image(path: Path, width: int | None, height: int | None, slide_sources: list[int]) -> tuple[str, str, str]:
    file_type = path.suffix.lower().lstrip(".")
    name = path.name.lower()
    area = (width or 0) * (height or 0)
    aspect = (width / height) if width and height else None

    if file_type == "emf":
        return "diagram", "redraw", "Legacy EMF/vector drawing; safer to rebuild as native editable geometry."
    if file_type == "wdp":
        return "unknown", "replace", "JPEG XR/HDPhoto asset is not portable for the rebuilt deck; keep as source evidence only."
    if area and area < 45_000:
        if "image1" in name or "image2" in name or "image3" in name:
            return "logo", "keep", "Small identity/brand asset; reuse only if visually clean."
        return "icon", "redraw", "Small low-resolution raster; redraw as native SVG or PowerPoint geometry."
    if aspect and aspect > 2.6:
        return "background", "enhance", "Wide image suitable for cinematic treatment after crop and contrast pass."
    if slide_sources and min(slide_sources) in {7, 8, 9} and area >= 250_000:
        return "photo", "enhance", "Lab visual from source deck; crop away legacy text overlay where needed and grade consistently."
    if slide_sources and min(slide_sources) in {11, 12} and area >= 100_000:
        return "photo", "enhance", "Product portfolio visual from source deck; preserve if it stays sharp at presentation scale."
    if file_type in {"jpg", "jpeg"}:
        return "photo", "enhance", "Photographic source; crop, darken, and grade consistently before reuse."
    if area >= 350_000:
        if slide_sources and min(slide_sources) in {7, 8, 9, 11, 12}:
            return "photo", "enhance", "Likely lab or product visual from source deck; inspect and use where relevant."
        return "screenshot", "redraw", "Large raster source that likely represents a diagram or screenshot; rebuild if text-heavy."
    return "decorative", "redraw", "PowerPoint-era decorative graphic; replace with the new technical visual system."


def extract_media() -> list[dict]:
    EXTRACTED.mkdir(parents=True, exist_ok=True)
    PROCESSED.mkdir(parents=True, exist_ok=True)
    media_entries = []
    with zipfile.ZipFile(SOURCE) as zf:
        media_by_slide = collect_slide_media(zf)
        slide_by_media = defaultdict(list)
        for slide_num, media_names in media_by_slide.items():
            for media_name in media_names:
                slide_by_media[media_name].append(slide_num)

        media_paths = sorted(
            [p for p in zf.namelist() if p.startswith("ppt/media/")],
            key=lambda p: Path(p).name,
        )
        for idx, media_path in enumerate(media_paths, 1):
            media_name = Path(media_path).name
            out_name = f"asset_{idx:03d}_{media_name}"
            out_path = EXTRACTED / out_name
            with zf.open(media_path) as src, out_path.open("wb") as dst:
                shutil.copyfileobj(src, dst)

            width = height = None
            image_mode = None
            has_alpha = False
            try:
                with Image.open(out_path) as im:
                    width, height = im.size
                    image_mode = im.mode
                    has_alpha = im.mode in ("RGBA", "LA") or ("transparency" in im.info)
            except Exception:
                pass

            file_type = out_path.suffix.lower().lstrip(".")
            slides = sorted(slide_by_media.get(media_name, []))
            classification, suggested_reuse, notes = classify_image(out_path, width, height, slides)
            aspect = round(width / height, 3) if width and height else None
            media_entries.append(
                {
                    "asset_id": f"asset_{idx:03d}",
                    "original_filename": media_name,
                    "extracted_filename": out_name,
                    "path": str(out_path.relative_to(ROOT)),
                    "slide_number": slides[0] if slides else None,
                    "slide_numbers": slides,
                    "dimensions": {"width": width, "height": height},
                    "aspect_ratio": aspect,
                    "file_type": file_type,
                    "image_mode": image_mode,
                    "has_alpha": has_alpha,
                    "estimated_resolution_quality": estimate_image_quality(width, height, file_type),
                    "classification": classification,
                    "suggested_reuse": suggested_reuse,
                    "notes": notes,
                }
            )
    return media_entries


def write_contact_sheet(entries: list[dict]) -> None:
    pdf_path = ASSETS / "asset-contact-sheet.pdf"
    page_w, page_h = landscape(A4)
    margin = 10 * mm
    cols = 4
    rows = 3
    gap = 6 * mm
    cell_w = (page_w - 2 * margin - (cols - 1) * gap) / cols
    cell_h = (page_h - 2 * margin - (rows - 1) * gap) / rows
    c = canvas.Canvas(str(pdf_path), pagesize=(page_w, page_h))
    c.setTitle("UE Labs asset contact sheet")
    for n, entry in enumerate(entries):
        pos = n % (cols * rows)
        if n and pos == 0:
            c.showPage()
        row = pos // cols
        col = pos % cols
        x = margin + col * (cell_w + gap)
        y = page_h - margin - (row + 1) * cell_h - row * gap

        c.setStrokeColor(colors.HexColor("#D6D6D6"))
        c.setFillColor(colors.HexColor("#F5F5F5"))
        c.rect(x, y, cell_w, cell_h, stroke=1, fill=1)

        img_path = ROOT / entry["path"]
        thumb_box_h = cell_h * 0.62
        try:
            with Image.open(img_path) as im:
                im = ImageOps.exif_transpose(im)
                im.thumbnail((int(cell_w - 8 * mm), int(thumb_box_h - 6 * mm)))
                temp = PROCESSED / f"contact_{entry['asset_id']}.png"
                bg = Image.new("RGB", im.size, (245, 245, 245))
                if im.mode in ("RGBA", "LA"):
                    bg.paste(im, mask=im.getchannel("A"))
                else:
                    bg.paste(im.convert("RGB"))
                bg.save(temp)
                iw, ih = im.size
                ix = x + (cell_w - iw) / 2
                iy = y + cell_h - thumb_box_h + (thumb_box_h - ih) / 2
                c.drawImage(str(temp), ix, iy, width=iw, height=ih, preserveAspectRatio=True, mask="auto")
        except Exception:
            c.setFillColor(colors.HexColor("#D9DEE3"))
            c.rect(x + 4 * mm, y + cell_h - thumb_box_h + 3 * mm, cell_w - 8 * mm, thumb_box_h - 6 * mm, fill=1, stroke=0)
            c.setFillColor(colors.HexColor("#4B5563"))
            c.setFont("Helvetica-Bold", 10)
            c.drawCentredString(x + cell_w / 2, y + cell_h - thumb_box_h / 2, entry["file_type"].upper())

        c.setFillColor(colors.HexColor("#1F2937"))
        c.setFont("Helvetica-Bold", 8.5)
        c.drawString(x + 3 * mm, y + 25 * mm, f"{entry['asset_id']} | slide {entry.get('slide_number') or 'n/a'}")
        c.setFont("Helvetica", 7.5)
        dims = entry["dimensions"]
        dim_text = f"{dims.get('width') or '?'} x {dims.get('height') or '?'} px"
        c.drawString(x + 3 * mm, y + 20 * mm, f"{entry['classification']} | {dim_text}")
        c.drawString(x + 3 * mm, y + 15 * mm, f"Reuse: {entry['suggested_reuse']}")
        c.setFont("Helvetica", 6.7)
        note = entry["notes"][:92]
        c.drawString(x + 3 * mm, y + 10 * mm, note)
    c.save()


def shape_texts(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = clean_text(shape.text)
            if text:
                texts.append(text)
    return texts


def slide_title(slide, texts: list[str]) -> str:
    if any("UE LABS" in t for t in texts):
        return "UE LABS"
    if slide.shapes.title is not None:
        title = clean_text(slide.shapes.title.text)
        if title:
            return title
    for text in texts:
        if len(text) <= 90 and not text.startswith("©") and "Confidential Property" not in text and not text.startswith("Page "):
            return text.split("\n")[0]
    return ""


def extract_content(media_entries: list[dict]) -> dict:
    prs = Presentation(str(SOURCE))
    with zipfile.ZipFile(SOURCE) as zf:
        notes = collect_notes(zf)

    assets_by_slide = defaultdict(list)
    for entry in media_entries:
        for slide_num in entry.get("slide_numbers") or []:
            assets_by_slide[slide_num].append(entry["asset_id"])

    all_texts = []
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = shape_texts(slide)
        all_texts.extend(texts)
        joined = "\n".join(texts)
        title = slide_title(slide, texts)
        body = [t for t in texts if t != title]
        numbers = sorted(set(re.findall(r"\b(?:\d{1,4}(?:[.,]\d+)?|DA307|C22\.2|94\.[12]|746)\b", joined)))
        labs = sorted({name for name in LAB_NAMES if re.search(re.escape(name), joined, re.I)})
        products = sorted({name for name in PRODUCT_NAMES if re.search(re.escape(name), joined, re.I)})
        standards = sorted({std for std in STANDARDS if std.replace("Nº", "N").lower() in joined.replace("Nº", "N").lower() or std.lower() in joined.lower()})
        if "UL 50E" in joined and "94.2" in joined:
            standards.append("UL 50E / CSA C22.2 Nº 94.2")
            standards = sorted(set(standards))
        footer_markings = [
            t
            for t in texts
            if "Confidential Property" in t or "All Rights Reserved" in t or re.fullmatch(r"Page \d+", t)
        ]
        accreditation_refs = [
            line
            for text in texts
            for line in clean_text(text).split("\n")
            if any(token in line for token in ["UL", "CTDP", "DA307", "Underwriters Laboratories"])
        ]
        slides.append(
            {
                "slide_number": idx,
                "title": title,
                "body_text": body,
                "speaker_script_text": notes.get(idx, ""),
                "visible_numbers": numbers,
                "lab_names": labs,
                "standards": standards,
                "product_names": products,
                "accreditation_references": accreditation_refs,
                "footer_confidentiality_markings": footer_markings,
                "images_assets_used": assets_by_slide.get(idx, []),
                "content_transformation_recommendation": recommend_transformation(idx, title, joined),
            }
        )

    normalized_counts = Counter(clean_text(t) for t in all_texts)
    duplicated = sorted([text for text, count in normalized_counts.items() if count > 1 and len(text) > 3])

    return {
        "source_file": SOURCE.name,
        "slide_count": len(prs.slides),
        "slides": slides,
        "duplicated_template_text": duplicated,
        "obsolete_placeholder_text": ["Slide 13 contains no visible source text; rebuild as intentional closing slide."],
        "repeated_mandatory_text": [
            "Confidential Property of Schneider Electric",
            "© 2020 Schneider Electric, All Rights Reserved",
            "Page markers",
        ],
        "content_to_move_to_appendix_or_notes": [
            "Long standards descriptions can be preserved in speaker notes while the slide uses a grouped capability matrix.",
        ],
        "content_to_transform_into_visual_diagrams": [
            "SE organization chart",
            "Universal Enclosures governance model",
            "UE Labs mission bullets",
            "UE Labs three-site network",
            "Standards scope",
            "Product families",
        ],
    }


def recommend_transformation(slide_num: int, title: str, joined: str) -> str:
    if slide_num in {3, 4, 5}:
        return "Rebuild as simplified hierarchy/operating-model diagram."
    if slide_num == 6:
        return "Convert mission bullets into value-chain visual plus lab network callout."
    if slide_num in {7, 8, 9}:
        return "Use lab-specific hero visual treatment with consistent fact badges."
    if slide_num == 10:
        return "Convert dense standards text to grouped capability matrix; preserve full descriptions in notes."
    if slide_num in {11, 12}:
        return "Convert product text into product portfolio/environment protection visuals."
    if slide_num == 13:
        return "Replace blank slide with designed closing slide."
    return "Preserve facts while redesigning layout and hierarchy from scratch."


def main() -> None:
    ASSETS.mkdir(exist_ok=True)
    REPORTS.mkdir(exist_ok=True)
    media_entries = extract_media()
    (ASSETS / "asset-inventory.json").write_text(json.dumps(media_entries, indent=2, ensure_ascii=False), encoding="utf-8")
    write_contact_sheet(media_entries)
    content = extract_content(media_entries)
    (REPORTS / "content-inventory.json").write_text(json.dumps(content, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Extracted {len(media_entries)} media assets")
    print(f"Wrote {ASSETS / 'asset-inventory.json'}")
    print(f"Wrote {ASSETS / 'asset-contact-sheet.pdf'}")
    print(f"Wrote {REPORTS / 'content-inventory.json'}")


if __name__ == "__main__":
    main()
